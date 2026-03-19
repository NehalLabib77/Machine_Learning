from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Presentation object will be created per theme
prs = None

# Default color scheme (will be overridden by apply_theme)
primary_blue = RGBColor(0, 123, 255)
neutral_gray = RGBColor(52, 58, 64)
light_bg = RGBColor(248, 249, 250)
green_accent = RGBColor(40, 167, 69)
red_accent = RGBColor(220, 53, 69)

# Apply a named theme (sets color variables). Returns title/text colors for convenience.
def apply_theme(name='modern_light'):
    global primary_blue, neutral_gray, light_bg, green_accent, red_accent
    if name == 'modern_light':
        primary_blue = RGBColor(0, 123, 255)
        neutral_gray = RGBColor(52, 58, 64)
        light_bg = RGBColor(248, 249, 250)
        green_accent = RGBColor(40, 167, 69)
        red_accent = RGBColor(220, 53, 69)
        title_color = primary_blue
        text_color = neutral_gray
    elif name == 'dark':
        primary_blue = RGBColor(88, 101, 242)
        neutral_gray = RGBColor(230, 230, 230)
        light_bg = RGBColor(18, 18, 20)
        green_accent = RGBColor(72, 195, 148)
        red_accent = RGBColor(255, 99, 132)
        title_color = RGBColor(255,255,255)
        text_color = RGBColor(235,235,235)
    elif name == 'warm':
        primary_blue = RGBColor(255, 99, 71)
        neutral_gray = RGBColor(60, 48, 50)
        light_bg = RGBColor(255, 250, 244)
        green_accent = RGBColor(255, 165, 79)
        red_accent = RGBColor(220, 53, 69)
        title_color = primary_blue
        text_color = neutral_gray
    else:
        # fallback
        return apply_theme('modern_light')
    return {'title_color': title_color, 'text_color': text_color}

# Build the presentation using the active theme
def build_presentation(theme='modern_light', out_path='stylish_ticket_booking_ppt_themed.pptx', copy_dest=None):
    global prs
    apply_theme(theme)
    prs = Presentation()
    # Now add slides (calls to helper functions below will use the current theme colors)
    
    # Slide 1: Title
    add_title_slide("Online Ticket Booking Mobile App", "Flutter-Based Software Project Proposal")
    
    # Slide 2: Motivation
    add_bullet_slide("Motivation", [
        "Skip the queues: Booking tickets is still tedious and unreliable.",
        "Users deserve a seamless, mobile-first experience.",
        "Our app simplifies it all—search, book, pay on the go."
    ])

    # Slide 3: Problem Statement
    add_bullet_slide("Problem Statement", [
        "Manual systems: Long waits and errors.",
        "No real-time updates: Seats vanish without warning.",
        "Slow payments: Frustrating and insecure."
    ], accent_color=red_accent)

    # Slide 4: Existing Systems & Gaps
    add_bullet_slide("Existing Systems & Gaps", [
        "Overloaded apps: Too many features, confusing UX.",
        "Web-heavy: Not optimized for mobile speed.",
        "Usability gap: Ignores on-the-go users."
    ])

    # Slide 5: Project Objective
    add_bullet_slide("Project Objective", [
        "Build a sleek Flutter app for effortless booking.",
        "Enable digital search, selection, and payment.",
        "Prioritize security and user delight."
    ], accent_color=green_accent)

    # Slide 6: Project Scope
    add_bullet_slide("Project Scope", [
        "User side: Register, authenticate, book, view history.",
        "Booking flow: Search tickets, select seats, pay online.",
        "Admin tools: Manage events and inventory."
    ])

    # Slide 7: Proposed Solution
    add_bullet_slide("Proposed Solution", [
        "Cross-platform: Flutter for iOS/Android.",
        "Versatile: Bus, movies, events.",
        "Integrated payments: Secure test mode."
    ])

    # Slide 8: Core Features
    add_bullet_slide("Core Features", [
        "Secure login/signup.",
        "Intuitive search & seat picker.",
        "Fast payment integration.",
        "Instant confirmations & history."
    ])

    # Slide 9: System Architecture (two-column modern layout)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Left column: bullets
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(4.2), Inches(4.5))
    ltf = left_box.text_frame
    ltf.clear()
    ltf.text = "System Architecture"
    ltf.paragraphs[0].font.size = Pt(26)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = neutral_gray
    ltf.paragraphs[0].font.name = 'Segoe UI'

    left_body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.2), Inches(3.3))
    ltf2 = left_body.text_frame
    ltf2.clear()
    for t in ["Frontend: Flutter App","Backend: Firebase / REST API","Database: Cloud Firestore","Payments: Stripe / Razorpay"]:
        p = ltf2.add_paragraph()
        p.text = f"• {t}"
        p.font.size = Pt(16)
        p.font.color.rgb = neutral_gray
        p.font.name = 'Segoe UI'

    # Right column: simple diagram represented by rounded rectangles
    x = Inches(5.0)
    y = Inches(0.9)
    w = Inches(3.2)
    h = Inches(0.8)
    services = ["Mobile App", "API / Backend", "Cloud DB", "Payments"]
    for i, s in enumerate(services):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + i * Inches(1.05), w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = light_bg
        shp.line.color.rgb = primary_blue
        tx = shp.text_frame
        tx.clear()
        tx.text = s
        tx.paragraphs[0].font.size = Pt(12)
        tx.paragraphs[0].font.bold = True
        tx.paragraphs[0].font.color.rgb = primary_blue

    add_footer(slide)

    # Slide 10: Technology Stack
    add_bullet_slide("Technology Stack", [
        "App: Flutter & Dart.",
        "Auth: Firebase.",
        "DB: Cloud Firestore.",
        "Payments: Gateway SDK."
    ])

    # Slide 11: Development Methodology
    add_bullet_slide("Development Methodology", [
        "Agile: Flexible sprints.",
        "Iterative: Build, test, refine.",
        "Focus: Continuous improvement."
    ])

    # Slide 12: Implementation Timeline
    add_bullet_slide("Implementation Timeline", [
        "Month 1: Plan & Design.",
        "Month 2: Develop Core.",
        "Month 3: Integrate Backend.",
        "Month 4: Test & Polish."
    ])

    # Slide 13: Testing Strategy
    add_bullet_slide("Testing Strategy", [
        "Unit: Code logic checks.",
        "Integration: Firebase flows.",
        "User: Real feedback loops."
    ])

    # Slide 14: Risk Management
    add_bullet_slide("Risk Management", [
        "Payments: Sandbox testing.",
        "Time: Prioritize essentials.",
        "Skills: Hands-on Flutter practice."
    ])

    # Slide 15: Expected Outcomes
    add_bullet_slide("Expected Outcomes", [
        "Working app prototype.",
        "Real-world skills in payments/tech.",
        "Portfolio boost."
    ], accent_color=green_accent)

    # Slide 16: Conclusion
    add_bullet_slide("Conclusion", [
        "Solves real booking pains.",
        "Feasible in 4 months.",
        "Ready to launch—approve?"
    ])

    # Slide 17: Thank You
    add_title_slide("Thank You", "Questions?\nContact: your_email@example.com")

    # Save the themed presentation
    prs.save(out_path)
    print('Themed PPT generated:', out_path)

    # Optionally copy to destination folder
    if copy_dest:
        import shutil
        try:
            shutil.copyfile(out_path, copy_dest)
            print('Copied themed PPT to:', copy_dest)
        except Exception as e:
            print('Copy failed:', e)

    return prs

# Helper: add a small footer to a slide
def add_footer(slide):
    footer = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.6), prs.slide_width - Inches(1), Inches(0.4))
    tf = footer.text_frame
    tf.clear()
    tf.text = "Nehal Labib (ID: XXXXX) | Supervisor: XXXXX | CSE 499 – Project / Thesis"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = neutral_gray
    p.alignment = PP_ALIGN.CENTER

# Helper: set a neutral slide background (modern light tone)
def set_slide_bg(slide, color=light_bg):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

# Helper function to add a modern title slide
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout for full control
    set_slide_bg(slide)

    # Top accent band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = primary_blue
    try:
        band.line.fill.background()
    except Exception:
        pass

    # Title on top band (white)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), prs.slide_width - Inches(1.2), Inches(0.6))
    ttf = title_box.text_frame
    ttf.clear()
    ttf.text = title
    p = ttf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.name = 'Segoe UI'

    # Subtitle under the band
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), prs.slide_width - Inches(1.6), Inches(1))
    stf = subtitle_box.text_frame
    stf.clear()
    stf.text = subtitle
    stf.paragraphs[0].font.size = Pt(16)
    stf.paragraphs[0].font.italic = True
    stf.paragraphs[0].font.color.rgb = neutral_gray
    stf.paragraphs[0].font.name = 'Segoe UI'

    add_footer(slide)

# Helper function to add a modern bullet slide
def add_bullet_slide(title, bullets, accent_color=primary_blue):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide)

    # Title area
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(0.8))
    ttf = title_box.text_frame
    ttf.clear()
    ttf.text = title
    p = ttf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = neutral_gray
    p.font.name = 'Segoe UI'

    # Accent bar under title
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(1.6), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    try:
        accent.line.fill.background()
    except Exception:
        pass

    # Body area
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), prs.slide_width - Inches(1.6), Inches(4))
    btf = body_box.text_frame
    btf.clear()  # prevent the common first-empty-paragraph issue
    for bullet in bullets:
        p = btf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = neutral_gray
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)

    add_footer(slide)

# Slides are now created inside build_presentation(theme='...') to ensure the presentation object exists.
# Call build_presentation(...) from the main guard at the bottom of this file.

# (Slides are now created inside build_presentation)

if __name__ == '__main__':
    # Create a modern light-themed PPT and copy a fixed version to your D: folder
    build_presentation(theme='modern_light', out_path='stylish_ticket_booking_ppt_themed.pptx', copy_dest=r'D:\5th assignment\New folder\stylish_ticket_booking_ppt_themed.pptx')
